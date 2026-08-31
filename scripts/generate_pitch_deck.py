"""Generate the editable CrossSignal hackathon pitch deck.

Run ``pip install python-pptx`` once before invoking this script.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submission" / "CrossSignal-Hackathon-Pitch-Final.pptx"
COVER = ROOT / "assets" / "crosssignal-hackathon-cover.png"

NAVY = RGBColor(7, 29, 73)
BLUE = RGBColor(0, 59, 112)
CYAN = RGBColor(25, 181, 216)
LIGHT_CYAN = RGBColor(114, 212, 232)
ICE = RGBColor(242, 247, 251)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(83, 101, 125)
GREEN = RGBColor(30, 142, 83)
AMBER = RGBColor(198, 126, 0)


def rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=22, color=NAVY, bold=False,
         font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def label(slide, value, number):
    text(slide, f"{number:02d}", .55, .36, .5, .25, 10, CYAN, True)
    text(slide, value.upper(), 1.05, .34, 5.5, .3, 10, MUTED, True)
    rect(slide, .55, .75, 12.2, .015, RGBColor(213, 224, 233))


def title(slide, headline, subhead=None):
    text(slide, headline, .65, 1.08, 12, 1.05, 34, NAVY, True, "Aptos Display")
    if subhead:
        text(slide, subhead, .68, 2.02, 11.7, .55, 15, MUTED)


def bullet_list(slide, items, x, y, w, h, size=19, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(14)
        p.text = f"•  {item}"
    return box


def footer(slide, number):
    text(slide, "CROSSSIGNAL · OMOBOLAJI E ADEYAN · ALPACA PAPER TRADING", .65, 7.08,
         10.8, .2, 8, MUTED, True)
    text(slide, str(number), 12.1, 7.04, .5, .25, 9, MUTED, True, align=PP_ALIGN.RIGHT)


def add_notes(slide, notes):
    frame = slide.notes_slide.notes_text_frame
    frame.text = notes


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Hook
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, .7, .7, .06, .42, CYAN)
    text(s, "DECISION INTELLIGENCE FOR CROSS-MARKET TRADING", .9, .73, 7.2, .35, 11, LIGHT_CYAN, True)
    text(s, "Markets disagree.\nWe trade the gap.", .72, 1.55, 8.8, 2.25, 42, WHITE, True, "Aptos Display")
    text(s, "An auditable AI trading agent that challenges, governs and scores every decision.",
         .75, 4.18, 8.2, .75, 20, LIGHT_CYAN)
    rect(s, 9.7, -.6, 4.5, 4.5, CYAN, radius=True)
    rect(s, 10.42, .12, 3.1, 3.1, NAVY, radius=True)
    text(s, "CROSSSIGNAL", .75, 6.25, 4, .4, 18, WHITE, True)
    text(s, "Created by Omobolaji E Adeyan", .75, 6.7, 5, .3, 12, LIGHT_CYAN)
    add_notes(s, "Markets often disagree before they reprice. CrossSignal identifies that disagreement, challenges the conclusion, applies deterministic controls, and records whether the prediction was ultimately correct.")

    # 2 — Problem
    s = prs.slides.add_slide(blank); label(s, "The problem", 2)
    title(s, "Trading agents generate signals.\nFew can defend a decision.")
    cards = [
        ("Fragmented", "Single-market analysis misses information moving between rates, credit and equities."),
        ("Unaccountable", "LLM conclusions are difficult to reproduce, challenge or audit after the outcome."),
        ("Unmeasured", "Most agents do not precommit a forecast and later score whether it was right."),
    ]
    for i, (head, body) in enumerate(cards):
        x = .65 + i * 4.15
        rect(s, x, 3.25, 3.8, 2.55, ICE, RGBColor(213, 224, 233))
        text(s, f"0{i+1}", x + .25, 3.55, .5, .35, 12, CYAN, True)
        text(s, head, x + .25, 4.0, 3.2, .4, 21, NAVY, True)
        text(s, body, x + .25, 4.55, 3.2, .85, 14, MUTED)
    footer(s, 2)
    add_notes(s, "Markets transmit information at different speeds. CrossSignal addresses fragmentation, weak AI accountability, and the absence of honest forward scoring.")

    # 3 — Solution
    s = prs.slides.add_slide(blank); label(s, "The solution", 3)
    title(s, "A scientific proof chain for every trade", "The SIGNAL protocol turns a thesis into a governed, falsifiable decision.")
    stages = [("Source", "Provenance"), ("Compare", "Disagreement"), ("Challenge", "Falsification"),
              ("Stress", "Greeks"), ("Seal", "SHA-256"), ("Score", "Outcome")]
    for i, (head, sub) in enumerate(stages):
        x = .55 + i * 2.1
        rect(s, x, 3.05, 1.72, 1.5, NAVY if i in (2, 4) else ICE,
             NAVY if i in (2, 4) else RGBColor(213, 224, 233))
        text(s, f"{i+1}", x + .18, 3.25, .3, .3, 11, CYAN, True)
        text(s, head, x + .18, 3.67, 1.35, .3, 16, WHITE if i in (2, 4) else NAVY, True)
        text(s, sub, x + .18, 4.05, 1.35, .25, 9, LIGHT_CYAN if i in (2, 4) else MUTED)
        if i < 5:
            text(s, "→", x + 1.76, 3.55, .3, .3, 16, CYAN, True)
    rect(s, .65, 5.25, 12.0, .72, BLUE)
    text(s, "ABSTAIN is a successful outcome when evidence or execution quality is insufficient.",
         .95, 5.47, 11.4, .3, 16, WHITE, True, align=PP_ALIGN.CENTER)
    footer(s, 3)
    add_notes(s, "Every decision passes a proof chain. Crucially, abstention is treated as a correct policy outcome when evidence fails a gate.")

    # 4 — Alpaca architecture
    s = prs.slides.add_slide(blank); label(s, "Technology", 4)
    title(s, "Alpaca is the market and execution backbone")
    rect(s, .65, 2.45, 3.25, 3.35, ICE, RGBColor(213, 224, 233))
    text(s, "LIVE EVIDENCE", .95, 2.8, 2.7, .3, 11, CYAN, True)
    bullet_list(s, ["Stocks and options", "Quotes, volume and Greeks", "Account and positions", "Alpaca News context"], .95, 3.35, 2.65, 1.8, 15)
    rect(s, 5.0, 2.45, 3.25, 3.35, NAVY)
    text(s, "CROSSSIGNAL", 5.3, 2.8, 2.65, .3, 11, LIGHT_CYAN, True)
    bullet_list(s, ["Claude macro synthesis", "Deterministic construction", "Risk and stability gates", "Sealed Decision Contract"], 5.3, 3.35, 2.65, 1.8, 15, WHITE)
    rect(s, 9.35, 2.45, 3.25, 3.35, ICE, RGBColor(213, 224, 233))
    text(s, "PAPER LIFECYCLE", 9.65, 2.8, 2.65, .3, 11, CYAN, True)
    bullet_list(s, ["All-leg preflight", "Multi-leg paper orders", "Fill reconciliation", "Explicit recovery state"], 9.65, 3.35, 2.65, 1.8, 15)
    text(s, "→", 4.18, 3.8, .5, .5, 25, CYAN, True)
    text(s, "→", 8.53, 3.8, .5, .5, 25, CYAN, True)
    footer(s, 4)
    add_notes(s, "The official Alpaca MCP server powers sponsor-native market evidence and the paper execution lifecycle. Claude reasons, but deterministic code controls broker access.")

    # 5 — Differentiation/evidence
    s = prs.slides.add_slide(blank); label(s, "Defensibility", 5)
    title(s, "One scorecard. Two defensible decisions.",
          "Signal strength never gets to hide weak execution evidence.")
    metrics = [("85", "Signal quality"), ("100", "Decision stability"),
               ("93", "Execution quality"), ("Pending", "Outcome evidence")]
    for i, (value, caption) in enumerate(metrics):
        x = .65 + i * 3.02
        rect(s, x, 2.65, 2.72, 1.3, ICE, RGBColor(213, 224, 233))
        text(s, value, x + .2, 2.88, 2.32, .45, 24, NAVY, True,
             align=PP_ALIGN.CENTER)
        text(s, caption, x + .2, 3.42, 2.32, .25, 11, MUTED, True,
             align=PP_ALIGN.CENTER)
    rect(s, .65, 4.35, 5.85, 1.38, NAVY)
    text(s, "AUTHORIZED REPLAY", .95, 4.62, 2.2, .25, 11, LIGHT_CYAN, True)
    text(s, "All deterministic gates passed", .95, 5.03, 4.9, .3, 17, WHITE, True)
    rect(s, 6.78, 4.35, 5.85, 1.38, ICE, RGBColor(213, 224, 233))
    text(s, "LIVE CLOUD EVIDENCE", 7.08, 4.62, 2.3, .25, 11, CYAN, True)
    text(s, "ABSTAIN · 14 of 15 gates passed", 7.08, 5.03, 4.95, .3, 17, NAVY, True)
    text(s, "CS-20260831-66AAE940 · sealed receipt · zero broker mutations",
         .75, 6.18, 11.7, .25, 11, MUTED, True, align=PP_ALIGN.CENTER)
    footer(s, 5)
    add_notes(s, "This is the key proof slide. The same policy evaluates both cases. The authorized replay shows the complete path when every gate passes. The live GitHub Evidence Watch contract scored 85 for signal and 100 for stability, but execution quality was 93 because one of fifteen gates failed—so the correct autonomous decision was abstention. Outcome evidence remains pending until the sealed horizon matures.")

    # 6 — Safety
    s = prs.slides.add_slide(blank); label(s, "Trust and safety", 6)
    title(s, "Safe by architecture—not by prompt")
    controls = [("Public UI", "Read-only replay", GREEN), ("Cloud watch", "Observe only", GREEN),
                ("LLM", "No direct tool access", GREEN), ("Evidence", "Secret-free receipts", GREEN),
                ("Recovery", "Explicit approval", AMBER), ("News", "Untrusted context", AMBER)]
    for i, (head, body, tone) in enumerate(controls):
        row, col = divmod(i, 3)
        x, y = .65 + col * 4.12, 2.6 + row * 1.45
        rect(s, x, y, 3.75, 1.12, ICE, RGBColor(213, 224, 233))
        rect(s, x, y, .08, 1.12, tone)
        text(s, head, x + .3, y + .2, 1.25, .28, 11, NAVY, True)
        text(s, body, x + 1.65, y + .2, 1.75, .5, 11, MUTED)
    rect(s, .65, 5.67, 12.0, .58, BLUE)
    text(s, "NIST AI RMF and SSDF aligned · voluntary and tailorable · not NIST-certified",
         .9, 5.85, 11.5, .25, 13, WHITE, True, align=PP_ALIGN.CENTER)
    footer(s, 6)
    add_notes(s, "CrossSignal follows NIST-aligned risk-management reasoning without claiming certification. Public infrastructure and scheduled GitHub automation have no broker mutation capability. The cloud watcher always runs execute false and exports only a secret-free artifact.")

    # 7 — Business
    s = prs.slides.add_slide(blank); label(s, "Business value", 7)
    title(s, "The accountability layer for AI-assisted trading")
    text(s, "WHO NEEDS IT", .7, 2.5, 3.5, .3, 11, CYAN, True)
    bullet_list(s, ["Active options traders", "Small investment teams", "Broker and platform providers", "Risk and compliance teams"], .7, 3.0, 4.6, 2.2, 18)
    rect(s, 5.35, 2.45, .02, 3.4, RGBColor(213, 224, 233))
    text(s, "HOW IT GROWS", 5.8, 2.5, 3.5, .3, 11, CYAN, True)
    steps = [("01", "Professional SaaS"), ("02", "Team governance"),
             ("03", "Broker partnerships"), ("04", "Enterprise policy APIs")]
    for i, (num, value) in enumerate(steps):
        y = 3.0 + i * .67
        text(s, num, 5.8, y, .65, .3, 11, CYAN, True)
        text(s, value, 6.55, y - .03, 4.35, .35, 17, NAVY, True)
    rect(s, 10.55, 2.45, 2.05, 3.4, NAVY)
    text(s, "VALUE", 10.9, 2.8, 1.35, .3, 10, LIGHT_CYAN, True, align=PP_ALIGN.CENTER)
    text(s, "Faster\nreview.\nStronger\ntrust.", 10.72, 3.45, 1.7, 1.75, 18, WHITE, True, align=PP_ALIGN.CENTER)
    footer(s, 7)
    add_notes(s, "The wedge is decision accountability for serious individual and small-team options users, with a path toward broker and enterprise governance integrations.")

    # 8 — Proof and CTA
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    text(s, "THE PROOF IS PUBLIC", .75, .7, 4, .3, 11, LIGHT_CYAN, True)
    text(s, "Inspect the intelligence.\nChallenge the decision.", .72, 1.35, 8.8, 1.55, 36, WHITE, True, "Aptos Display")
    stats = [("38", "automated tests"), ("85/100", "live signal quality"), ("14/15", "risk gates passed")]
    for i, (value, caption) in enumerate(stats):
        x = .75 + i * 2.55
        text(s, value, x, 3.35, 2.2, .55, 27, LIGHT_CYAN, True)
        text(s, caption, x, 3.95, 2.2, .3, 12, WHITE)
    rect(s, 8.85, 1.2, 3.75, 4.6, WHITE)
    text(s, "LIVE JUDGE DEMO", 9.2, 1.62, 3.05, .3, 11, BLUE, True, align=PP_ALIGN.CENTER)
    text(s, "crosssignal-ai-agent\n.streamlit.app", 9.2, 2.25, 3.05, .85, 18, NAVY, True, align=PP_ALIGN.CENTER)
    rect(s, 9.2, 3.25, 3.05, .03, CYAN)
    text(s, "SOURCE CODE", 9.2, 3.72, 3.05, .3, 11, BLUE, True, align=PP_ALIGN.CENTER)
    text(s, "github.com/omobolajiadeyan/\nalpaca-cross-market-agent", 9.05, 4.28, 3.35, .85, 13, NAVY, True, align=PP_ALIGN.CENTER)
    text(s, "CrossSignal does not merely recommend a trade.\nIt proves whether the trade deserves authorization.",
         .75, 5.35, 7.45, .85, 17, WHITE, True)
    text(s, "Omobolaji E Adeyan · Paper trading only · Not investment advice", .75, 6.75, 8, .3, 11, LIGHT_CYAN)
    add_notes(s, "Close on verified evidence: thirty-eight tests pass in GitHub, the scheduled cloud watcher sealed a live contract, and the system abstained because one risk gate failed. CrossSignal makes AI trading decisions inspectable, challengeable, reproducible, and measurable. Invite judges to open the live application, repository, and Actions evidence artifact.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
