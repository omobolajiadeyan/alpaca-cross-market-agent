"""Generate the one-page CrossSignal technical write-up as an editable PPTX.

The portrait, single-slide PPTX is exported to PDF during package assembly.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submission" / "CrossSignal-One-Page-Writeup.pptx"
NAVY = RGBColor(7, 29, 73)
CYAN = RGBColor(25, 181, 216)
MUTED = RGBColor(69, 84, 105)
PALE = RGBColor(239, 247, 250)
WHITE = RGBColor(255, 255, 255)


def box(slide, text, x, y, w, h, size=9.2, color=MUTED, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(0)
    return shape


def heading(slide, text, x, y, w):
    box(slide, text.upper(), x, y, w, 0.22, 10.5, NAVY, True)
    line = slide.shapes.add_shape(1, Inches(x), Inches(y + 0.28), Inches(w), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.color.rgb = CYAN


def build():
    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE

    banner = slide.shapes.add_shape(1, 0, 0, Inches(8.5), Inches(1.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.color.rgb = NAVY
    box(slide, "CROSSSIGNAL", 0.45, 0.28, 3.2, 0.38, 23, WHITE, True)
    box(slide, "An auditable cross-market options agent", 0.47, 0.72, 4.4, 0.25, 11, RGBColor(174, 226, 238), False)
    box(slide, "Omobolaji E Adeyan", 6.0, 0.32, 2.0, 0.25, 10.5, WHITE, True)
    box(slide, "Alpaca paper: PA3PDTUDIXDU", 5.35, 0.68, 2.65, 0.22, 9, RGBColor(174, 226, 238), False)

    left_x, right_x, col_w = 0.45, 4.4, 3.65
    heading(slide, "Strategy and AI logic", left_x, 1.45, col_w)
    box(slide,
        "CrossSignal asks whether a promising signal deserves execution. It synchronizes six lenses: equity implied volatility, Treasury curve, credit, realized volatility, rate expectations, and options positioning. A deterministic engine ranks cross-market disagreements. Claude proposes a structured thesis; a separate adversarial pass returns the strongest counterargument, missing evidence, alternative explanation, invalidation condition, and confidence penalty. Claude never receives broker authority.\n\n"
        "The SIGNAL protocol applies ten bounded perturbations, then seals the market snapshot, thesis, challenge, stability result, proposed portfolio, prediction horizon, and invalidation rule into a SHA-256 Decision Contract before any broker call. Later scoring compares the precommitted direction with inverse and cash counterfactuals.",
        left_x, 1.82, col_w, 3.0)

    heading(slide, "Risk gates and options", left_x, 4.95, col_w)
    box(slide,
        "Signals map only to defined-risk SPY, HYG, and TLT vertical spreads. Base gates require complete structure, maximum proposed loss ≤ $1,500, post-challenge confidence ≥ 55%, sufficient buying power, diversification, and zero fallback feeds. All legs then pass quote, Greeks, liquidity, bid-ask, stress, margin, and drawdown checks. One failed gate produces ABSTAIN. After entry, a persisted policy takes 50% of maximum profit, stops at 50% of maximum loss, exits after five trading days, or closes two days before expiry. Broker-leg reconciliation, fresh timestamps, an atomic claim, and a deterministic client order ID guard each multi-leg close.",
        left_x, 5.32, col_w, 2.35, 9.0)

    heading(slide, "Alpaca infrastructure", left_x, 7.95, col_w)
    box(slide,
        "Alpaca's official MCP server supplies market data, options Greeks, account state, paper orders, spread valuation, the market clock, and reconciliation. Mutations require the exact paper endpoint and explicit local authorization; automated exits require a second switch. A separate entry kill switch preserves management of open positions. The public demo and GitHub Evidence Watch are read-only.",
        left_x, 8.32, col_w, 1.65)

    panel = slide.shapes.add_shape(1, Inches(4.18), Inches(1.28), Inches(3.87), Inches(8.75))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALE
    panel.line.color.rgb = PALE
    heading(slide, "September 3 evidence", right_x, 1.45, col_w)
    box(slide,
        "LOCAL CONTRACT — CS-20260903-FE01A097\n"
        "Signal 82.4/100 · stability 90% · all feeds sourced live\n"
        "Confidence 68% → 53% · five of six base gates\n\n"
        "The adversarial review rejected the ‘cheap vol’ framing: IV rank came from a short local history, the put/call value was a narrow ATM proxy, and IV was already 1.26× realized volatility. Confidence fell below the fixed 55% floor, so no option preflight or order occurred.\n\n"
        "CLOUD CONTRACT — CS-20260903-5C194F65\n"
        "Signal 82 · stability 90 · execution 87 · 13/15 checks\n"
        "Confidence cleared at 56% and Greeks coverage was 6/6. Preflight rejected displayed volume 0 < 10 and maximum relative spread 93.33% > 25%. The after-hours option market was not executable; no order was queued.",
        right_x, 1.82, col_w, 4.15, 9.05)

    heading(slide, "Honest limitation", right_x, 6.35, col_w)
    box(slide,
        "The fresh competition account has $100,000 cash, zero positions, and zero orders. CrossSignal therefore has no competition-account P&L to claim. This weakens the P&L judging dimension but is not an explicit eligibility failure. The agent did not loosen a gate or manufacture an after-hours paper trade for the demo.",
        right_x, 6.72, col_w, 1.55, 9.2)

    heading(slide, "Links", right_x, 8.55, col_w)
    box(slide,
        "Demo: crosssignal-ai-agent.streamlit.app\n"
        "Code: github.com/omobolajiadeyan/alpaca-cross-market-agent\n"
        "Evidence: GitHub Actions run 33805502192\n\n"
        "Paper trading only. Simulated results do not guarantee future performance. Not investment advice.",
        right_x, 8.92, col_w, 1.25, 8.8)

    box(slide, "59 automated tests · governed entry and exit · $100,000 fresh Alpaca paper account", 0.45, 10.45, 7.6, 0.22, 9.4, NAVY, True)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
